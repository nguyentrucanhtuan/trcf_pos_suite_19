#!/usr/bin/env python3
"""
Convert Markdown slides to PPTX format
"""

import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image

def parse_markdown_slides(md_content):
    """Parse markdown content into slides"""
    # Split by slide separator (---)
    slides_raw = re.split(r'\n---\n', md_content)
    slides = []
    
    for slide_raw in slides_raw:
        slide_raw = slide_raw.strip()
        if not slide_raw:
            continue
            
        slide = {
            'title': '',
            'image': None,
            'content': [],
            'script': ''
        }
        
        lines = slide_raw.split('\n')
        current_section = None
        
        for line in lines:
            # Parse title (## Slide X: Title)
            if line.startswith('## '):
                slide['title'] = re.sub(r'^## (Slide \d+: )?', '', line).strip()
            # Parse image
            elif line.startswith('!['):
                match = re.search(r'\!\[.*?\]\((.*?)\)', line)
                if match:
                    slide['image'] = match.group(1)
            # Parse sections
            elif line.startswith('**Script:**') or line.startswith('> "'):
                current_section = 'script'
                if line.startswith('> "'):
                    slide['script'] += line[3:].strip() + ' '
            elif line.startswith('> ') and current_section == 'script':
                slide['script'] += line[2:].strip() + ' '
            elif line.startswith('**Nội dung chính:**'):
                current_section = 'content'
            elif current_section == 'content' and line.startswith('- '):
                slide['content'].append(line[2:].strip())
            elif current_section == 'content' and line.startswith('| '):
                slide['content'].append(line.strip())
                
        if slide['title']:
            slides.append(slide)
    
    return slides

def create_pptx(slides, output_path, images_dir):
    """Create PowerPoint presentation from slides"""
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # Get blank layout
    blank_layout = prs.slide_layouts[6]  # Blank
    
    for slide_data in slides:
        slide = prs.slides.add_slide(blank_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data['title']
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0x1a, 0x36, 0x5d)  # Dark blue
        title_para.alignment = PP_ALIGN.LEFT
        
        # Add image if exists
        if slide_data['image']:
            img_path = slide_data['image']
            # Resolve relative path from slides directory (parent of images)
            slides_dir = os.path.dirname(images_dir)
            if img_path.startswith('./'):
                img_path = os.path.join(slides_dir, img_path[2:])
            elif not os.path.isabs(img_path):
                img_path = os.path.join(slides_dir, img_path)
                
            if os.path.exists(img_path):
                try:
                    # Get image dimensions
                    with Image.open(img_path) as img:
                        img_width, img_height = img.size
                    
                    # Calculate size to fit
                    max_width = Inches(10)
                    max_height = Inches(5.5)
                    
                    aspect = img_width / img_height
                    if aspect > (10/5.5):
                        width = max_width
                        height = width / aspect
                    else:
                        height = max_height
                        width = height * aspect
                    
                    # Center horizontally
                    left = (prs.slide_width - width) / 2
                    slide.shapes.add_picture(img_path, left, Inches(1.2), width=width, height=height)
                except Exception as e:
                    print(f"Warning: Could not add image {img_path}: {e}")
        
        # Add content bullets
        if slide_data['content']:
            content_top = Inches(7) if slide_data['image'] else Inches(1.5)
            content_box = slide.shapes.add_textbox(Inches(0.5), content_top, Inches(15), Inches(1.5))
            tf = content_box.text_frame
            tf.word_wrap = True
            
            for i, item in enumerate(slide_data['content'][:5]):  # Max 5 items
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = "• " + item
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        
        # Add script as notes
        if slide_data['script']:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data['script'].strip('"').strip()
    
    prs.save(output_path)
    print(f"Created: {output_path}")
    return output_path

def main():
    md_file = '/Users/tuan/coffeetree_odoo19_dev/custom_addons/docs/slides/them_san_pham/slides.md'
    images_dir = '/Users/tuan/coffeetree_odoo19_dev/custom_addons/docs/slides/them_san_pham/images'
    output_file = '/Users/tuan/coffeetree_odoo19_dev/custom_addons/docs/slides/them_san_pham/slides.pptx'
    
    print("Reading markdown...")
    with open(md_file, 'r', encoding='utf-8') as f:
        md_content = f.read()
    
    print("Parsing slides...")
    slides = parse_markdown_slides(md_content)
    print(f"Found {len(slides)} slides")
    
    print("Creating PPTX...")
    create_pptx(slides, output_file, images_dir)
    
    print(f"\nDone! PPTX file created at:\n{output_file}")

if __name__ == '__main__':
    main()
